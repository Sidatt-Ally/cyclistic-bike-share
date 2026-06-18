"""
export_pptx.py — Présentation PowerPoint Cyclistic Bike-Share
Auteur : Sidi Mohamed ALLY | Analyste de données Professionnel Certifié (Google)
Sortie  : outputs/Cyclistic_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Chemins ──────────────────────────────────────────────────────────────────
BASE   = r"C:\Users\Republic Of Computer\Desktop\Cyclistic_Project"
FIGS   = os.path.join(BASE, "figures")
OUT    = os.path.join(BASE, "outputs")
os.makedirs(OUT, exist_ok=True)

BADGE  = os.path.join(FIGS, "badge_google_da.png")
FIG    = {i: os.path.join(FIGS, f) for i, f in {
    1: "01_rides_count.png",
    2: "02_avg_ride_length.png",
    3: "03_rides_by_dow.png",
    4: "04_rides_by_month.png",
    5: "05_bike_types.png",
    6: "06_top10_stations.png",
    7: "07_rides_by_hour.png",
    8: "08_avg_length_by_dow.png",
}.items()}

# ── Palette ───────────────────────────────────────────────────────────────────
BLUE        = RGBColor(0x1A, 0x73, 0xE8)   # membre
ORANGE      = RGBColor(0xE8, 0x4D, 0x0E)   # casual
DARK        = RGBColor(0x1A, 0x1A, 0x2E)   # fond couverture
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY    = RGBColor(0x75, 0x75, 0x75)
ACCENT      = RGBColor(0x0D, 0x47, 0xA1)   # bleu foncé titres

# ── Dimensions (widescreen 16:9) ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout complètement vide

# ═════════════════════════════════════════════════════════════════════════════
# Helpers
# ═════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    """Ajoute un rectangle coloré."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    """Ajoute une zone de texte."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def heading(slide, text, y=0.22, size=28, color=ACCENT):
    """Titre de slide standard."""
    txbox(slide, text, 0.4, y, 12.5, 0.6, size=size, bold=True, color=color)

def divider(slide, y=0.95, color=BLUE):
    """Ligne horizontale de séparation sous le titre."""
    rect(slide, 0.4, y, 12.5, 0.04, fill=color)

def figure(slide, fig_path, x, y, w):
    """Insère une image et calcule la hauteur proportionnelle."""
    from PIL import Image as PILImage
    if not os.path.exists(fig_path):
        print(f"[WARN] Figure introuvable : {fig_path}")
        return
    with PILImage.open(fig_path) as im:
        iw, ih = im.size
    h = w * ih / iw
    slide.shapes.add_picture(fig_path, Inches(x), Inches(y), Inches(w), Inches(h))

def notes(slide, text):
    """Ajoute des notes d'orateur (invisibles au public)."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = text

def kpi_box(slide, x, y, value, label, val_color=BLUE):
    """Boîte KPI : grande valeur + petit label."""
    rect(slide, x, y, 2.8, 1.3, fill=LIGHT_GRAY)
    txbox(slide, value, x+0.1, y+0.1, 2.6, 0.7,
          size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txbox(slide, label, x+0.1, y+0.78, 2.6, 0.45,
          size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

def bullet(slide, items, x, y, w=11, size=16, color=DARK, spacing=0.42):
    """Liste à puces simple."""
    cy = y
    for item in items:
        txbox(slide, f"▸  {item}", x, cy, w, 0.38, size=size, color=color)
        cy += spacing

def two_col_table(slide, headers, rows, x, y, col_widths, row_h=0.42):
    """Tableau simple sans dépendance à pptx.Table (évite les conflits de layout)."""
    # En-têtes
    cx = x
    for i, (h_txt, cw) in enumerate(zip(headers, col_widths)):
        rect(slide, cx, y, cw-0.02, row_h, fill=BLUE)
        txbox(slide, h_txt, cx+0.08, y+0.06, cw-0.18, row_h-0.1,
              size=13, bold=True, color=WHITE)
        cx += cw
    # Lignes
    for ri, row in enumerate(rows):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        cx = x
        cy = y + row_h * (ri + 1)
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            rect(slide, cx, cy, cw-0.02, row_h, fill=bg)
            c = BLUE if ci == 0 else DARK
            txbox(slide, cell, cx+0.08, cy+0.06, cw-0.18, row_h-0.1,
                  size=12, color=c, bold=(ci == 0))
            cx += cw


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Couverture
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=DARK)                      # fond sombre

# Bandeau bleu gauche
rect(s, 0, 0, 0.25, 7.5, fill=BLUE)

# Badge Google DA
if os.path.exists(BADGE):
    s.shapes.add_picture(BADGE, Inches(10.8), Inches(0.3), Inches(2.2), Inches(2.2))

# Titre principal
txbox(s, "Cyclistic Bike-Share", 0.6, 1.5, 9.5, 1.0,
      size=44, bold=True, color=WHITE)
txbox(s, "Comment convertir les cyclistes occasionnels\nen membres annuels ?",
      0.6, 2.55, 9.5, 1.1, size=22, color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)

# Ligne décorative
rect(s, 0.6, 3.75, 5.0, 0.06, fill=BLUE)

# Auteur & certification
txbox(s, "Sidi Mohamed ALLY", 0.6, 3.95, 8, 0.5,
      size=18, bold=True, color=WHITE)
txbox(s, "Analyste de données  |  Google Data Analytics Certificate",
      0.6, 4.48, 9, 0.4, size=14, color=RGBColor(0xBB, 0xDE, 0xFB))

# KPIs couverture
for i, (val, lbl, col) in enumerate([
    ("3,88 M", "trajets analysés", WHITE),
    ("12 mois", "2019 – 2020",     WHITE),
    ("2 groupes", "membre / casual", BLUE),
]):
    bx = 0.6 + i * 4.1
    rect(s, bx, 5.5, 3.8, 1.3, fill=RGBColor(0x0D, 0x1B, 0x3E))
    txbox(s, val, bx+0.1, 5.55, 3.6, 0.65,
          size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(s, lbl, bx+0.1, 6.2, 3.6, 0.45,
          size=12, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

notes(s, """Bonjour à tous.
Je m'appelle Sidi Mohamed ALLY, analyste de données certifié Google.
Aujourd'hui je vous présente mon étude de cas sur Cyclistic Bike-Share — la capstone du Google Data Analytics Certificate.

La question centrale : comment transformer les cyclistes occasionnels en membres annuels ?
Pour y répondre, j'ai analysé 3,88 millions de trajets sur 12 mois (2019-2020).

Les deux groupes étudiés : les membres (abonnés annuels) et les casual (tickets à l'unité ou à la journée).""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Contexte business
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Contexte business", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

txbox(s, "Cyclistic — Chicago, Illinois", 0.4, 1.25, 12, 0.45,
      size=20, bold=True, color=ACCENT)

bullet(s, [
    "5 824 vélos répartis sur 692 stations dans Chicago",
    "Deux types d'utilisateurs : membres annuels et casual (ticket jour / trajet unique)",
    "Les membres génèrent plus de revenus — la direction veut augmenter leur part",
    "Lily Moreno (Directrice Marketing) : convertir les casual en membres = levier de croissance prioritaire",
], x=0.5, y=1.78, size=16)

# Schéma simple casual → membre
rect(s, 1.0, 5.0, 3.5, 1.1, fill=RGBColor(0xFF, 0xE0, 0xCC))
txbox(s, "Casual", 1.0, 5.05, 3.5, 0.5, size=22, bold=True,
      color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s, "ticket jour / trajet", 1.0, 5.5, 3.5, 0.45, size=12,
      color=MID_GRAY, align=PP_ALIGN.CENTER)

txbox(s, "➜  Campagne\n     marketing ciblée", 4.6, 5.15, 3.5, 0.85,
      size=15, color=ACCENT, align=PP_ALIGN.CENTER)

rect(s, 8.3, 5.0, 3.5, 1.1, fill=RGBColor(0xE3, 0xF2, 0xFD))
txbox(s, "Membre", 8.3, 5.05, 3.5, 0.5, size=22, bold=True,
      color=BLUE, align=PP_ALIGN.CENTER)
txbox(s, "abonnement annuel", 8.3, 5.5, 3.5, 0.45, size=12,
      color=MID_GRAY, align=PP_ALIGN.CENTER)

notes(s, """Cyclistic est un service de vélo en libre-service fictif basé à Chicago.
Le modèle économique repose sur deux types d'utilisateurs.

Les membres paient un abonnement annuel — ils sont la base stable de revenus.
Les casual paient à l'usage — plus coûteux à acquérir, moins rentables sur la durée.

La directrice marketing, Lily Moreno, a identifié que la conversion des casual en membres
est le levier de croissance le plus efficace plutôt que d'aller chercher de nouveaux clients.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Question analytique & méthodologie
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Question analytique & Méthodologie", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

# Question centrale
rect(s, 0.4, 1.25, 12.5, 0.95, fill=RGBColor(0xE3, 0xF2, 0xFD))
txbox(s,
      "« Comment les membres annuels et les cyclistes occasionnels utilisent-ils\n"
      "  les vélos Cyclistic différemment ? »",
      0.55, 1.32, 12.2, 0.82, size=17, bold=True, color=ACCENT)

# Méthodologie en 4 étapes
steps = [
    ("1. Collecter",   "11 fichiers CSV\n3,9 M trajets bruts"),
    ("2. Nettoyer",    "Python / pandas\nSQL PostgreSQL"),
    ("3. Analyser",    "Statistiques descriptives\nVisualisation matplotlib"),
    ("4. Partager",    "Rapport DOCX + PDF\nTableau Public"),
]
for i, (titre, desc) in enumerate(steps):
    bx = 0.4 + i * 3.2
    col = [BLUE, RGBColor(0x0D,0x47,0xA1), RGBColor(0x15,0x65,0xC0), RGBColor(0x19,0x76,0xD2)][i]
    rect(s, bx, 2.45, 3.0, 1.5, fill=col)
    txbox(s, titre, bx+0.1, 2.52, 2.8, 0.5,
          size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, desc, bx+0.1, 3.0, 2.8, 0.85,
          size=12, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

# Outils
txbox(s, "Outils utilisés :", 0.4, 4.2, 3, 0.4, size=14, bold=True, color=ACCENT)
tools_row = ["Python 3 / pandas", "PostgreSQL / SQL", "matplotlib / seaborn",
             "python-docx / ReportLab", "Tableau Public"]
cx = 0.4
for tool in tools_row:
    w = len(tool) * 0.13 + 0.4
    rect(s, cx, 4.65, w, 0.42, fill=LIGHT_GRAY, line=BLUE, line_w=Pt(1))
    txbox(s, tool, cx+0.1, 4.7, w-0.15, 0.32, size=11, color=ACCENT)
    cx += w + 0.18

notes(s, """La question analytique est au cœur de l'étude. Elle est de type descriptif-comparatif :
on cherche des différences de comportement entre deux groupes.

La méthodologie suit le cadre du cours Google DA :
Ask → Prepare → Process → Analyze → Share → Act

Pour le nettoyage : suppression des trajets < 1 min (tests de maintenance) et > 24h,
correction des types de données, ajout des colonnes calculées (ride_length, day_of_week, saison).

Résultat net après nettoyage : 3,885 millions de trajets valides sur 3,906 millions bruts.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Résultat 1 : Volume global
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 1 — Volume : qui roule le plus ?", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

# KPIs
kpi_box(s, 0.4,  1.3, "2,29 M",  "trajets membres  (59%)", val_color=BLUE)
kpi_box(s, 3.35, 1.3, "1,59 M",  "trajets casual   (41%)", val_color=ORANGE)
kpi_box(s, 6.3,  1.3, "3,88 M",  "total trajets",           val_color=ACCENT)

# Figure
figure(s, FIG[1], x=0.4, y=2.8, w=8.5)

# Insight
rect(s, 9.5, 2.8, 3.6, 2.2, fill=RGBColor(0xE8, 0xF4, 0xFD))
rect(s, 9.5, 2.8, 0.08, 2.2, fill=BLUE)
txbox(s, "Insight", 9.7, 2.88, 3.3, 0.4, size=13, bold=True, color=BLUE)
txbox(s,
      "Les membres représentent 59% du volume — "
      "ils sont la colonne vertébrale du service.\n\n"
      "Les casual (41%) constituent un bassin de\nconversion significatif.",
      9.7, 3.3, 3.3, 1.6, size=12, color=DARK)

notes(s, """Premier constat : les membres génèrent plus de trajets que les casual — 59% vs 41%.

Ce ratio est important : il signifie qu'il y a déjà une base solide de membres,
mais que 4 trajets sur 10 sont faits par des casual qui pourraient potentiellement
devenir membres.

1,59 million de trajets casual sur l'année — c'est un bassin de conversion non négligeable.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Résultat 2 : Durée des trajets
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 2 — Durée : les casual roulent plus longtemps", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

kpi_box(s, 0.4,  1.3, "23,1 min", "durée moyenne — casual",  val_color=ORANGE)
kpi_box(s, 3.35, 1.3, "15,6 min", "durée moyenne — membres", val_color=BLUE)
kpi_box(s, 6.3,  1.3, "× 1,48",   "ratio casual / membre",   val_color=ACCENT)

figure(s, FIG[2], x=0.4, y=2.8, w=8.5)

rect(s, 9.5, 2.8, 3.6, 2.5, fill=RGBColor(0xFF, 0xF3, 0xE0))
rect(s, 9.5, 2.8, 0.08, 2.5, fill=ORANGE)
txbox(s, "Interprétation", 9.7, 2.88, 3.3, 0.4, size=13, bold=True, color=ORANGE)
txbox(s,
      "Les casual utilisent le vélo pour\nle loisir (trajets longs, détendus).\n\n"
      "Les membres l'utilisent pour les\ntrajets domicile-travail (courts,\nefficaces).\n\n"
      "Comportement fondamentalement\ndifférent → message marketing\ndifférent.",
      9.7, 3.3, 3.3, 1.9, size=12, color=DARK)

notes(s, """Résultat clé : un casual roule en moyenne 23,1 minutes contre 15,6 minutes pour un membre.
Ratio : 1,48 — les casual passent presque 50% plus de temps sur le vélo.

Ce n'est pas une question de distance, c'est une question de comportement :
- Les membres : trajet domicile → bureau → maison. Court, régulier, prévisible.
- Les casual : balade, tourisme, exploration. Plus long, moins fréquent.

Cette différence fondamentale doit guider le message marketing :
pour un casual, l'abonnement annuel doit être présenté comme un accès illimité au loisir,
pas comme un outil de transport quotidien.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Résultat 3 : Jour de la semaine
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 3 — Temporel : semaine vs week-end", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

figure(s, FIG[3], x=0.4, y=1.25, w=7.8)
figure(s, FIG[8], x=0.4, y=4.3,  w=7.8)

# Tableau comparatif
two_col_table(s,
    headers=["Groupe", "Pic semaine", "Pic week-end", "Usage dominant"],
    rows=[
        ["Membres", "Jeudi (23.1%)", "Samedi",        "Domicile-travail"],
        ["Casual",  "Vendredi",      "Samedi (23.1%)", "Loisir / tourisme"],
    ],
    x=8.5, y=1.5,
    col_widths=[1.8, 1.8, 1.8, 1.8],
    row_h=0.52
)

rect(s, 8.5, 3.85, 4.6, 2.2, fill=RGBColor(0xE8, 0xF5, 0xE9))
rect(s, 8.5, 3.85, 0.08, 2.2, fill=RGBColor(0x2E, 0x7D, 0x32))
txbox(s, "Action marketing", 8.7, 3.93, 4.3, 0.4,
      size=13, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))
txbox(s,
      "Cibler les casual le week-end :\n"
      "• Offre découverte samedi-dimanche\n"
      "• Publicité dans les zones touristiques\n"
      "• Promotion d'abonnement week-end → annuel",
      8.7, 4.35, 4.3, 1.6, size=12, color=DARK)

notes(s, """Résultat très parlant sur le comportement hebdomadaire.

Les membres ont leur pic le jeudi — typique d'une utilisation domicile-travail
avec légère baisse le vendredi (télétravail, départ en week-end).

Les casual explosent le samedi — comportement de loisir pur.
Vendredi soir commence déjà leur pic.

Implication concrète : si on veut toucher des casual, il faut être présent le week-end.
Flyers dans les stations touristiques le samedi matin, offres limitées week-end,
notifications push samedi midi quand ils viennent de finir un trajet.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Résultat 4 : Saisonnalité
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 4 — Saisonnalité : l'été, saison clé", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

figure(s, FIG[4], x=0.4, y=1.25, w=9.0)

# Tableau saisons
two_col_table(s,
    headers=["Saison", "Membres", "Casual", "Écart"],
    rows=[
        ["Été (juin-août)",       "~42%", "~52%", "Casual ++ (loisir estival)"],
        ["Automne (sept-nov)",    "~32%", "~27%", "Membres résistants"],
        ["Printemps (mars-mai)",  "~16%", "~14%", "Équilibré"],
        ["Hiver (déc-fév)",       "~10%", "~7%",  "Chute des deux groupes"],
    ],
    x=9.6, y=1.35,
    col_widths=[2.2, 1.0, 1.0, 2.0],
    row_h=0.5
)

rect(s, 0.4, 5.6, 12.5, 1.3, fill=RGBColor(0xFF, 0xF8, 0xE1))
rect(s, 0.4, 5.6, 0.08, 1.3, fill=RGBColor(0xF5, 0x7F, 0x17))
txbox(s, "Fenêtre d'opportunité", 0.6, 5.68, 4, 0.4,
      size=13, bold=True, color=RGBColor(0xE6, 0x5C, 0x00))
txbox(s,
      "Juin, juillet, août : les casual sont les plus nombreux et les plus actifs."
      " C'est la meilleure période pour lancer une campagne de conversion"
      " avec offre spéciale abonnement annuel.",
      0.6, 6.1, 12.2, 0.65, size=13, color=DARK)

notes(s, """La saisonnalité est un facteur critique : l'été concentre la majorité des trajets casual.

En été, les casual représentent ~52% de leurs trajets annuels — c'est là qu'ils sont
les plus engagés avec le service, et donc les plus réceptifs à une offre d'abonnement.

Recommandation : lancer la campagne de conversion en mai pour capturer le pic de juin-août.
Ne pas attendre septembre — les casual disparaissent à l'automne et encore plus en hiver.

L'hiver : seulement 7% des trajets casual. Inutile d'investir en publicité à cette période.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Résultat 5 : Heure de départ
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 5 — Heure de départ : deux profils distincts", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

figure(s, FIG[7], x=0.4, y=1.25, w=8.8)

rect(s, 9.5, 1.35, 3.6, 5.5, fill=LIGHT_GRAY)
txbox(s, "Profil Membre", 9.65, 1.42, 3.3, 0.42,
      size=14, bold=True, color=BLUE)
bullet(s, [
    "Pic à 8h (rush matin)",
    "Pic à 17h (rush soir)",
    "Creux en journée",
    "→ trajet domicile-travail",
], x=9.65, y=1.88, w=3.2, size=12, spacing=0.38)

txbox(s, "Profil Casual", 9.65, 3.72, 3.3, 0.42,
      size=14, bold=True, color=ORANGE)
bullet(s, [
    "Montée progressive de 6h à 15h",
    "Pic unique à 15h-16h",
    "Pas de pic rush matin",
    "→ loisir / exploration",
], x=9.65, y=4.18, w=3.2, size=12, spacing=0.38)

notes(s, """Le graphique horaire est l'un des plus révélateurs de l'étude.

Les membres ont clairement deux pics : 8h et 17h — ce sont les heures de rush.
La signature d'un navetteur domicile-travail est parfaitement lisible.

Les casual ont un profil complètement différent : la courbe monte doucement de 6h à 15h,
avec un plateau l'après-midi. Aucun pic à 8h. C'est le profil d'une promenade de loisir.

Cette différence renforce la conclusion : casual et membres n'utilisent pas le même service
pour les mêmes raisons. Le message marketing doit le reconnaître.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Résultat 6 : Top stations casual
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 6 — Zones géographiques : où sont les casual ?", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

figure(s, FIG[6], x=0.4, y=1.25, w=8.2)

rect(s, 8.9, 1.35, 4.2, 5.6, fill=RGBColor(0xFF, 0xF3, 0xE0))
rect(s, 8.9, 1.35, 0.08, 5.6, fill=ORANGE)
txbox(s, "Ce que révèlent\nles stations", 9.1, 1.42, 3.8, 0.65,
      size=14, bold=True, color=ORANGE)
bullet(s, [
    "Streeter Dr & Grand Ave : bord du lac Michigan",
    "Millennium Park : zone touristique n°1",
    "Michigan Ave / Oak St : promenades shopping",
    "Shedd Aquarium, Navy Pier → loisir pur",
    "",
    "→ Les casual partent des zones\n   touristiques et des parcs",
    "→ Publicité physique ciblée\n   dans ces zones",
    "→ Partenariat offices de tourisme",
], x=9.1, y=2.15, w=3.8, size=11, spacing=0.38)

notes(s, """La carte des top 10 stations casual est très instructive.

Presque toutes les stations les plus fréquentées par les casual sont situées :
- En bord du lac Michigan (Streeter Dr, DuSable Lake Shore Dr)
- Dans les zones touristiques (Millennium Park, Navy Pier, Shedd Aquarium)
- Sur les grandes avenues commerçantes (Michigan Ave, Clark St)

Ce ne sont PAS les stations des quartiers résidentiels où habitent les membres.

Implication : la publicité physique (affiches, QR codes, flyers) doit être concentrée
dans ces zones touristiques, pas dans les quartiers de bureau.
Un partenariat avec les offices de tourisme de Chicago serait très pertinent.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Types de vélos
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Résultat 7 — Types de vélos : les casual aiment les docked bikes", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

figure(s, FIG[5], x=0.4, y=1.25, w=9.0)

rect(s, 9.6, 1.35, 3.5, 3.0, fill=LIGHT_GRAY)
txbox(s, "Observation", 9.75, 1.42, 3.2, 0.42, size=14, bold=True, color=ACCENT)
bullet(s, [
    "Les casual utilisent davantage\nles docked bikes",
    "Les membres préfèrent\nles classic bikes",
    "Les electric bikes sont\nappréciés des deux groupes",
], x=9.75, y=1.9, w=3.2, size=12, spacing=0.52)

notes(s, """Les types de vélos révèlent une préférence différente selon le groupe.

Les casual utilisent plus les docked bikes — vélos à station fixe, adaptés aux balades
sans contrainte de temps, souvent utilisés par des touristes qui ne connaissent pas bien la ville.

Les membres préfèrent les classic bikes — légers, efficaces pour les trajets quotidiens.

Les electric bikes gagnent en popularité chez les deux groupes, ce qui est une tendance
à surveiller pour les futures décisions d'investissement en flotte.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Recommandations
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 1.1, fill=ACCENT)
heading(s, "Recommandations — 3 actions prioritaires", y=0.22, size=26, color=WHITE)
divider(s, y=1.12, color=BLUE)

recs = [
    (
        "1. Campagne estivale de conversion",
        BLUE,
        [
            "Lancer en mai, cibler juin-août (pic casual)",
            "Offre : -20% sur l'abonnement annuel pour tout casual actif",
            "Canal : notification push après le 3e trajet du mois",
            "Message : « Vous aimez Cyclistic — passez à l'illimité »",
        ]
    ),
    (
        "2. Marketing géolocalisé week-end",
        ORANGE,
        [
            "Cibler les stations touristiques (Millennium Park, Navy Pier…)",
            "Affiches + QR code vers page d'abonnement",
            "Ambassadeurs le samedi matin dans les zones à fort trafic",
            "Partenariat offices de tourisme de Chicago",
        ]
    ),
    (
        "3. Offre week-end → abonnement annuel",
        RGBColor(0x2E, 0x7D, 0x32),
        [
            "Créer un pass week-end (sam+dim) à prix attractif",
            "Inclure une offre de bascule vers l'annuel après 3 week-ends",
            "Tracker la conversion : % pass week-end → membre annuel",
            "Évaluer après 6 mois et ajuster le prix du pass",
        ]
    ),
]

for i, (titre, color, points) in enumerate(recs):
    bx = 0.4 + i * 4.3
    rect(s, bx, 1.25, 4.0, 0.55, fill=color)
    txbox(s, titre, bx+0.1, 1.3, 3.8, 0.45,
          size=13, bold=True, color=WHITE)
    for j, pt in enumerate(points):
        cy = 1.9 + j * 0.42
        txbox(s, f"• {pt}", bx+0.1, cy, 3.85, 0.38, size=11, color=DARK)

# Matrice effort/impact
rect(s, 0.4, 5.5, 12.5, 1.65, fill=LIGHT_GRAY)
txbox(s, "Matrice Effort × Impact", 0.55, 5.58, 4, 0.4,
      size=13, bold=True, color=ACCENT)

matrix = [
    ("Action",              "Effort", "Impact", "Priorité"),
    ("Campagne estivale",   "Moyen",  "Élevé",  "⭐ Haute"),
    ("Marketing géolocalisé","Faible","Élevé",  "⭐ Haute"),
    ("Pass week-end",       "Élevé",  "Moyen",  "Moyenne"),
]
two_col_table(s,
    headers=matrix[0],
    rows=matrix[1:],
    x=0.4, y=5.97,
    col_widths=[4.5, 2.5, 2.5, 2.5],
    row_h=0.4
)

notes(s, """Trois recommandations concrètes, ancrées dans les données.

Recommandation 1 — Campagne estivale : c'est la plus efficace car elle cible le bon moment.
Les casual sont en pic en juin-août. Une offre de réduction à ce moment-là, déclenchée
après plusieurs trajets, a une forte probabilité de conversion.

Recommandation 2 — Marketing géolocalisé : effort faible, impact élevé.
Les stations sont connues, les zones sont identifiées. Il suffit d'y mettre les bons supports.

Recommandation 3 — Pass week-end : plus complexe à mettre en place car cela nécessite
de créer un nouveau produit tarifaire. Impact moyen mais intéressant sur le long terme
pour habituer les casual à une logique d'abonnement.

La matrice effort×impact confirme que les actions 1 et 2 sont à prioriser.""")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion & signature
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=DARK)
rect(s, 0, 0, 0.25, 7.5, fill=BLUE)

txbox(s, "Conclusion", 0.6, 0.5, 9, 0.65,
      size=32, bold=True, color=WHITE)
rect(s, 0.6, 1.22, 5.5, 0.06, fill=BLUE)

synthese = [
    "Les membres et les casual ont des comportements fondamentalement différents",
    "Casual → loisir, week-end, été, zones touristiques, trajets longs",
    "Membres → navetteurs, semaine, trajets courts et réguliers",
    "3 recommandations ciblées : campagne estivale, marketing géo, pass week-end",
    "Priorité : cibler les casual en été dans les zones à fort trafic",
]
cy = 1.4
for pt in synthese:
    txbox(s, f"✓  {pt}", 0.6, cy, 11.5, 0.42, size=14, color=WHITE)
    cy += 0.44

# Séparateur
rect(s, 0.6, 3.72, 11.5, 0.05, fill=RGBColor(0x1A, 0x73, 0xE8))

# Signature
txbox(s, "Sidi Mohamed ALLY", 0.6, 3.9, 9, 0.55,
      size=22, bold=True, color=WHITE)
txbox(s, "Analyste de données  |  Google Data Analytics Certificate",
      0.6, 4.48, 9, 0.4, size=14, color=RGBColor(0xBB, 0xDE, 0xFB))
txbox(s, "sidatt6969@gmail.com",
      0.6, 4.92, 5, 0.38, size=13, color=RGBColor(0xBB, 0xDE, 0xFB))

if os.path.exists(BADGE):
    s.shapes.add_picture(BADGE, Inches(10.8), Inches(3.8), Inches(2.2), Inches(2.2))

txbox(s, "Merci — Questions ?", 0.6, 6.3, 9, 0.7,
      size=26, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

notes(s, """Pour conclure :

L'analyse de 3,88 millions de trajets révèle deux profils d'utilisateurs clairement distincts.
Les casual sont des utilisateurs de loisir, actifs le week-end et en été.
Les membres sont des navetteurs, actifs en semaine toute l'année.

Cette différence est une opportunité : les casual connaissent déjà Cyclistic.
Ils aiment le service. Il faut leur montrer que l'abonnement annuel correspond à leur usage.

Les trois recommandations sont actionnables immédiatement avec un budget marketing maîtrisé.

Je suis disponible pour répondre à vos questions.
Merci de votre attention.""")


# ═════════════════════════════════════════════════════════════════════════════
# Sauvegarde
# ═════════════════════════════════════════════════════════════════════════════
out_path = os.path.join(OUT, "Cyclistic_Presentation.pptx")
prs.save(out_path)
size_kb = os.path.getsize(out_path) // 1024
print(f"OK — {out_path}")
print(f"     {size_kb} KB  |  {len(prs.slides)} slides")
print(f"\nSlides :")
titles = [
    "01. Couverture",
    "02. Contexte business",
    "03. Question analytique & Méthodologie",
    "04. Résultat 1 — Volume",
    "05. Résultat 2 — Durée",
    "06. Résultat 3 — Jour de la semaine",
    "07. Résultat 4 — Saisonnalité",
    "08. Résultat 5 — Heure de départ",
    "09. Résultat 6 — Top stations casual",
    "10. Résultat 7 — Types de vélos",
    "11. Recommandations",
    "12. Conclusion & Signature",
]
for t in titles:
    print(f"  {t}")
